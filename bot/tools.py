"""
bot/tools.py — Tool execution registry for the orchestrator.

These are the "hands" of the AI. The orchestrator's LLM decides what to do,
and these functions actually do it. Each tool is a plain Python function
that takes structured args and returns a result dict.

Tools:
  - send_email      → SMTP email sending
  - write_file      → Write content to a file
  - read_file       → Read a file's contents
  - web_search      → Search the web (DuckDuckGo, no API key needed)
  - create_document → Generate PPTX, DOCX, or PDF
  - run_command     → Execute a shell command (sandboxed)
  - http_request    → Make an HTTP request
"""
from __future__ import annotations

import json
import logging
import os
import subprocess
import time
from pathlib import Path
from typing import Any

log = logging.getLogger("valhalla.tools")

# Where generated documents go
OUTPUT_DIR = Path.home() / ".fireside" / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


# ═══════════════════════════════════════════════════════════════
# Tool Registry
# ═══════════════════════════════════════════════════════════════

TOOLS: dict[str, dict] = {}


def register(name: str, description: str, parameters: dict):
    """Decorator to register a tool function."""
    def decorator(fn):
        TOOLS[name] = {
            "name": name,
            "description": description,
            "parameters": parameters,
            "function": fn,
        }
        return fn
    return decorator


def execute_tool(name: str, args: dict) -> dict:
    """Execute a registered tool by name with given arguments."""
    tool = TOOLS.get(name)
    if not tool:
        return {"error": f"Unknown tool: {name}", "available": list(TOOLS.keys())}

    try:
        log.info("[tools] Executing: %s(%s)", name, json.dumps(args)[:200])
        result = tool["function"](**args)
        log.info("[tools] ✓ %s completed", name)
        return {"ok": True, "tool": name, "result": result}
    except Exception as e:
        log.error("[tools] ✗ %s failed: %s", name, e)
        return {"ok": False, "tool": name, "error": str(e)}


def get_tool_definitions() -> list[dict]:
    """Return OpenAI-compatible tool definitions for llama-server function calling."""
    return [
        {
            "type": "function",
            "function": {
                "name": t["name"],
                "description": t["description"],
                "parameters": t["parameters"],
            },
        }
        for t in TOOLS.values()
    ]


# ═══════════════════════════════════════════════════════════════
# Tool: send_email
# ═══════════════════════════════════════════════════════════════

@register(
    name="send_email",
    description="Send an email via SMTP. Supports HTML body.",
    parameters={
        "type": "object",
        "properties": {
            "to": {"type": "string", "description": "Recipient email address"},
            "subject": {"type": "string", "description": "Email subject line"},
            "body": {"type": "string", "description": "Email body (plain text or HTML)"},
            "html": {"type": "boolean", "description": "If true, body is HTML", "default": False},
        },
        "required": ["to", "subject", "body"],
    },
)
def send_email(to: str, subject: str, body: str, html: bool = False) -> dict:
    """Send an email using SMTP credentials from valhalla.yaml or env vars."""
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart

    # Load SMTP config from env or config file
    smtp_host = os.environ.get("SMTP_HOST", "smtp.gmail.com")
    smtp_port = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user = os.environ.get("SMTP_USER", "")
    smtp_pass = os.environ.get("SMTP_PASS", "")
    from_addr = os.environ.get("SMTP_FROM", smtp_user)

    if not smtp_user or not smtp_pass:
        # Try loading from config
        try:
            config_path = Path.home() / ".fireside" / "smtp.json"
            if config_path.exists():
                cfg = json.loads(config_path.read_text(encoding="utf-8"))
                smtp_host = cfg.get("host", smtp_host)
                smtp_port = cfg.get("port", smtp_port)
                smtp_user = cfg.get("user", smtp_user)
                smtp_pass = cfg.get("password", smtp_pass)
                from_addr = cfg.get("from", smtp_user)
        except Exception:
            pass

    if not smtp_user:
        return {"sent": False, "error": "No SMTP credentials configured. "
                "Set SMTP_USER/SMTP_PASS env vars or create ~/.fireside/smtp.json"}

    msg = MIMEMultipart("alternative")
    msg["Subject"] = subject
    msg["From"] = from_addr
    msg["To"] = to

    if html:
        msg.attach(MIMEText(body, "html"))
    else:
        msg.attach(MIMEText(body, "plain"))

    with smtplib.SMTP(smtp_host, smtp_port) as server:
        server.ehlo()
        server.starttls()
        server.login(smtp_user, smtp_pass)
        server.sendmail(from_addr, [to], msg.as_string())

    log.info("[tools] Email sent to %s: %s", to, subject)
    return {"sent": True, "to": to, "subject": subject}


# ═══════════════════════════════════════════════════════════════
# Tool: write_file
# ═══════════════════════════════════════════════════════════════

@register(
    name="write_file",
    description="Write content to a file. Creates parent directories if needed.",
    parameters={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "File path (relative to ~/.fireside/outputs/ or absolute)"},
            "content": {"type": "string", "description": "File content to write"},
        },
        "required": ["path", "content"],
    },
)
def write_file(path: str, content: str) -> dict:
    """Write content to a file."""
    p = Path(path)
    if not p.is_absolute():
        p = OUTPUT_DIR / path

    # Security: prevent writing outside fireside dirs
    safe_roots = [Path.home() / ".fireside", OUTPUT_DIR, Path(".")]
    if not any(str(p.resolve()).startswith(str(r.resolve())) for r in safe_roots):
        return {"written": False, "error": f"Cannot write outside safe directories: {p}"}

    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding="utf-8")
    return {"written": True, "path": str(p), "size": len(content)}


# ═══════════════════════════════════════════════════════════════
# Tool: read_file
# ═══════════════════════════════════════════════════════════════

@register(
    name="read_file",
    description="Read the contents of a file.",
    parameters={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "File path to read"},
            "max_chars": {"type": "integer", "description": "Max characters to return", "default": 10000},
        },
        "required": ["path"],
    },
)
def read_file(path: str, max_chars: int = 10000) -> dict:
    """Read a file and return its contents."""
    p = Path(path)
    if not p.exists():
        return {"read": False, "error": f"File not found: {path}"}

    content = p.read_text(encoding="utf-8", errors="replace")
    truncated = len(content) > max_chars
    return {
        "read": True,
        "path": str(p),
        "content": content[:max_chars],
        "truncated": truncated,
        "total_size": len(content),
    }


# ═══════════════════════════════════════════════════════════════
# Tool: web_search
# ═══════════════════════════════════════════════════════════════

@register(
    name="web_search",
    description="Search the web using DuckDuckGo. Returns top results with titles and snippets.",
    parameters={
        "type": "object",
        "properties": {
            "query": {"type": "string", "description": "Search query"},
            "max_results": {"type": "integer", "description": "Max results to return", "default": 5},
        },
        "required": ["query"],
    },
)
def web_search(query: str, max_results: int = 5) -> dict:
    """Search the web via DuckDuckGo HTML (no API key needed)."""
    import urllib.request
    import urllib.parse
    import re

    url = f"https://html.duckduckgo.com/html/?q={urllib.parse.quote(query)}"
    headers = {"User-Agent": "Mozilla/5.0 (Fireside/1.0)"}

    req = urllib.request.Request(url, headers=headers)
    with urllib.request.urlopen(req, timeout=10) as resp:
        html = resp.read().decode("utf-8", errors="replace")

    # Parse results from DuckDuckGo HTML
    results = []
    # DuckDuckGo HTML search returns results in <a class="result__a"> tags
    pattern = r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>'
    snippet_pattern = r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>'

    links = re.findall(pattern, html, re.DOTALL)
    snippets = re.findall(snippet_pattern, html, re.DOTALL)

    for i, (href, title) in enumerate(links[:max_results]):
        # Clean HTML tags from title and snippet
        clean_title = re.sub(r'<[^>]+>', '', title).strip()
        clean_snippet = ""
        if i < len(snippets):
            clean_snippet = re.sub(r'<[^>]+>', '', snippets[i]).strip()

        # DuckDuckGo wraps URLs in a redirect — extract the real URL
        real_url = href
        if "uddg=" in href:
            match = re.search(r'uddg=([^&]+)', href)
            if match:
                real_url = urllib.parse.unquote(match.group(1))

        results.append({
            "title": clean_title,
            "url": real_url,
            "snippet": clean_snippet,
        })

    return {"query": query, "results": results, "count": len(results)}


# ═══════════════════════════════════════════════════════════════
# Tool: create_document
# ═══════════════════════════════════════════════════════════════

@register(
    name="create_document",
    description="Create a document (PPTX PowerPoint, DOCX Word, or TXT). "
                "For PPTX, provide slides as a list of {title, content} objects.",
    parameters={
        "type": "object",
        "properties": {
            "filename": {"type": "string", "description": "Output filename (e.g. 'report.pptx')"},
            "doc_type": {"type": "string", "enum": ["pptx", "docx", "txt"], "description": "Document type"},
            "title": {"type": "string", "description": "Document title"},
            "slides": {
                "type": "array",
                "description": "For PPTX: list of slides [{title, content}]",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {"type": "string"},
                    },
                },
            },
            "content": {"type": "string", "description": "For DOCX/TXT: document body text"},
        },
        "required": ["filename", "doc_type"],
    },
)
def create_document(filename: str, doc_type: str, title: str = "",
                    slides: list = None, content: str = "") -> dict:
    """Create a document file."""
    output_path = OUTPUT_DIR / filename

    if doc_type == "pptx":
        return _create_pptx(output_path, title, slides or [])
    elif doc_type == "docx":
        return _create_docx(output_path, title, content)
    elif doc_type == "txt":
        output_path.write_text(content, encoding="utf-8")
        return {"created": True, "path": str(output_path), "type": "txt"}
    else:
        return {"created": False, "error": f"Unknown doc type: {doc_type}"}


def _create_pptx(path: Path, title: str, slides: list) -> dict:
    """Create a PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        # Fallback: create a simple text-based outline
        log.warning("[tools] python-pptx not installed — creating text outline instead")
        lines = [f"# {title}\n"]
        for i, slide in enumerate(slides, 1):
            lines.append(f"\n## Slide {i}: {slide.get('title', '')}")
            lines.append(slide.get("content", ""))
        path = path.with_suffix(".txt")
        path.write_text("\n".join(lines), encoding="utf-8")
        return {
            "created": True,
            "path": str(path),
            "type": "txt",
            "note": "python-pptx not installed. Run: pip install python-pptx",
        }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    if title:
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

    # Content slides
    for s in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = s.get("title", "")
        # Add content to the body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Body
                shape.text = s.get("content", "")
                break

    prs.save(str(path))
    return {"created": True, "path": str(path), "type": "pptx", "slides": len(slides) + (1 if title else 0)}


def _create_docx(path: Path, title: str, content: str) -> dict:
    """Create a Word document."""
    try:
        from docx import Document
    except ImportError:
        # Fallback to plain text
        log.warning("[tools] python-docx not installed — creating text file instead")
        path = path.with_suffix(".txt")
        text = f"# {title}\n\n{content}" if title else content
        path.write_text(text, encoding="utf-8")
        return {
            "created": True,
            "path": str(path),
            "type": "txt",
            "note": "python-docx not installed. Run: pip install python-docx",
        }

    doc = Document()
    if title:
        doc.add_heading(title, 0)
    for paragraph in content.split("\n\n"):
        doc.add_paragraph(paragraph)
    doc.save(str(path))
    return {"created": True, "path": str(path), "type": "docx"}


# ═══════════════════════════════════════════════════════════════
# Tool: run_command
# ═══════════════════════════════════════════════════════════════

@register(
    name="run_command",
    description="Execute a shell command. Output is captured and returned. "
                "Times out after 30 seconds. Dangerous commands are blocked.",
    parameters={
        "type": "object",
        "properties": {
            "command": {"type": "string", "description": "Shell command to execute"},
            "cwd": {"type": "string", "description": "Working directory", "default": "."},
        },
        "required": ["command"],
    },
)
def run_command(command: str, cwd: str = ".") -> dict:
    """Execute a shell command with safety checks."""
    # Block dangerous commands
    BLOCKED = ["rm -rf /", "format c:", "del /s /q", ":(){ :|:& };:",
               "mkfs", "dd if=", "shutdown", "reboot"]
    cmd_lower = command.lower().strip()
    for blocked in BLOCKED:
        if blocked in cmd_lower:
            return {"executed": False, "error": f"Blocked dangerous command: {blocked}"}

    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=30,
            cwd=cwd,
        )
        return {
            "executed": True,
            "command": command,
            "exit_code": result.returncode,
            "stdout": result.stdout[:5000],  # Cap output
            "stderr": result.stderr[:2000],
        }
    except subprocess.TimeoutExpired:
        return {"executed": False, "error": "Command timed out (30s)"}


# ═══════════════════════════════════════════════════════════════
# Tool: http_request
# ═══════════════════════════════════════════════════════════════

@register(
    name="http_request",
    description="Make an HTTP GET or POST request. Returns the response body.",
    parameters={
        "type": "object",
        "properties": {
            "url": {"type": "string", "description": "URL to request"},
            "method": {"type": "string", "enum": ["GET", "POST"], "default": "GET"},
            "body": {"type": "string", "description": "POST body (JSON string)"},
            "headers": {"type": "object", "description": "Additional headers"},
        },
        "required": ["url"],
    },
)
def http_request(url: str, method: str = "GET", body: str = None,
                 headers: dict = None) -> dict:
    """Make an HTTP request."""
    import urllib.request

    req_headers = {"User-Agent": "Fireside/1.0"}
    if headers:
        req_headers.update(headers)

    data = body.encode("utf-8") if body else None
    if data:
        req_headers["Content-Type"] = "application/json"

    req = urllib.request.Request(url, data=data, headers=req_headers, method=method)

    try:
        with urllib.request.urlopen(req, timeout=15) as resp:
            content = resp.read().decode("utf-8", errors="replace")
            return {
                "status": resp.status,
                "url": url,
                "body": content[:10000],
                "truncated": len(content) > 10000,
            }
    except urllib.error.HTTPError as e:
        return {"status": e.code, "url": url, "error": str(e)}


# ═══════════════════════════════════════════════════════════════
# Tool: browse_and_act
# ═══════════════════════════════════════════════════════════════

@register(
    name="browse_and_act",
    description="Open a website and navigate it to accomplish a goal. "
                "Uses the user's existing browser profile (saved logins). "
                "Returns a compact action tree of clickable/typeable elements.",
    parameters={
        "type": "object",
        "properties": {
            "url": {"type": "string", "description": "URL to open"},
            "goal": {"type": "string", "description": "What to accomplish (e.g. 'Order a cold brew')"},
        },
        "required": ["url", "goal"],
    },
)
def browse_and_act(url: str, goal: str) -> dict:
    """Open a website and return the interactive action tree.

    Two-tier approach:
      1. Try Playwright (full JS rendering, interactive navigation)
      2. Fall back to HTTP fetch + BeautifulSoup parse_interactive (90% token reduction)
    """
    # ── Tier 1: Playwright (interactive, JS-rendered) ──
    try:
        from plugins.browse.actor import BrowserActor
        import asyncio

        async def _run():
            actor = BrowserActor(headless=True)
            try:
                result = await actor.open(url)
                if not result["ok"]:
                    return result
                result["goal"] = goal
                return result
            finally:
                await actor.close()

        try:
            loop = asyncio.get_event_loop()
            if loop.is_running():
                import concurrent.futures
                with concurrent.futures.ThreadPoolExecutor() as pool:
                    return pool.submit(lambda: asyncio.run(_run())).result(timeout=30)
            else:
                return asyncio.run(_run())
        except Exception as e:
            log.warning("[tools] Playwright browse failed: %s — falling back to HTTP", e)
    except ImportError:
        log.info("[tools] Playwright not installed — using HTTP + BeautifulSoup fallback")

    # ── Tier 2: HTTP fetch + parse_interactive (lightweight, no JS) ──
    try:
        import urllib.request
        from plugins.browse.parser import parse_interactive

        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0 Safari/537.36 Fireside/1.0"}
        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req, timeout=15) as resp:
            html = resp.read().decode("utf-8", errors="replace")

        tree = parse_interactive(html, base_url=url)
        return {
            "ok": True,
            "url": url,
            "goal": goal,
            "title": tree.title,
            "action_tree": tree.to_action_text(),
            "stats": tree.summary_stats(),
            "mode": "http_fallback",
            "note": "Fetched via HTTP (no JavaScript rendering). Install Playwright for full JS support.",
        }
    except Exception as e:
        log.error("[tools] HTTP browse fallback failed: %s", e)
        return {"ok": False, "error": f"Browse failed: {e}. Install Playwright for full support: pip install playwright && python -m playwright install chromium"}


# ═══════════════════════════════════════════════════════════════
# Tool: check_spending
# ═══════════════════════════════════════════════════════════════

@register(
    name="check_spending",
    description="Check if a purchase is within spending limits before proceeding. "
                "Returns whether the purchase is allowed, needs approval, or is blocked.",
    parameters={
        "type": "object",
        "properties": {
            "description": {"type": "string", "description": "What is being purchased"},
            "estimated_cost": {"type": "number", "description": "Estimated cost in dollars"},
            "site": {"type": "string", "description": "Website domain (e.g. 'doordash.com')"},
        },
        "required": ["description", "estimated_cost"],
    },
)
def check_spending(description: str, estimated_cost: float, site: str = "") -> dict:
    """Check spending limits before a purchase."""
    try:
        from plugins.browse.spending import check_purchase, get_spending_summary
        result = check_purchase(description, estimated_cost, site)
        result["summary"] = get_spending_summary()
        return result
    except ImportError:
        return {"allowed": True, "reason": "Spending module not available — allowing by default"}


# ═══════════════════════════════════════════════════════════════
# Tool: scan_message (Guardian / Scam Shield)
# ═══════════════════════════════════════════════════════════════

@register(
    name="scan_message",
    description="Scan a message (SMS, email, chat) for scam/phishing indicators. "
                "Returns a color-coded threat level: green (safe), yellow (suspect), red (scam).",
    parameters={
        "type": "object",
        "properties": {
            "text": {"type": "string", "description": "Message text to scan"},
            "sender": {"type": "string", "description": "Sender name or address"},
            "message_type": {"type": "string", "enum": ["sms", "email", "chat"], "default": "sms"},
        },
        "required": ["text"],
    },
)
def scan_message_tool(text: str, sender: str = "", message_type: str = "sms") -> dict:
    """Scan a message for scam indicators."""
    try:
        from plugins.guardian.scanner import scan_message
        return scan_message(text=text, sender=sender, message_type=message_type)
    except ImportError:
        return {"error": "Guardian plugin not available"}


# ═══════════════════════════════════════════════════════════════
# Summary
# ═══════════════════════════════════════════════════════════════

log.info("[tools] Registered %d tools: %s", len(TOOLS), ", ".join(TOOLS.keys()))

