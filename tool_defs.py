"""
tool_defs.py — Shared tool definitions and executor.

Used by both the main chat agent (api/v1.py) and pipeline subagents
(plugins/pipeline/handler.py). Single source of truth for tool schemas
and execution logic.
"""
from __future__ import annotations

import json
import logging
import subprocess
import urllib.request
from pathlib import Path
from typing import Optional

log = logging.getLogger("valhalla.tools")


# ---------------------------------------------------------------------------
# Tool Schemas (OpenAI function-calling format)
# ---------------------------------------------------------------------------

TOOL_SCHEMAS = [
    {
        "type": "function",
        "function": {
            "name": "files_list",
            "description": "List files and folders. Use when the user says things like 'what's on my desktop', 'show me the project structure', 'what files are in this folder', 'show me what's in ~/projects'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Directory path to list, e.g. 'C:/Users/Jordan/Desktop'"},
                    "recursive": {"type": "boolean", "description": "Whether to list recursively", "default": False},
                    "pattern": {"type": "string", "description": "Glob pattern to filter, e.g. '*.py'", "default": "*"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "files_read",
            "description": "Read a file's contents. Use when the user says 'show me that file', 'read the README', 'what's in config.yaml', 'open the error log', 'check the package.json'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path to the file to read"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "files_write",
            "description": "Create or save a file. Use when the user says 'write me a script', 'make a file', 'save this code', 'create a config', 'generate a component'. Creates parent directories automatically.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path where the file should be saved"},
                    "content": {"type": "string", "description": "The content to write to the file"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "files_delete",
            "description": "Delete a file or folder. Use when the user says 'delete that file', 'remove the old logs', 'clean up the temp folder'. Requires explicit user confirmation — ask first, then call with confirmed=true.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path to the file to delete"},
                    "confirmed": {"type": "boolean", "description": "Set to true after user confirms deletion"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "terminal_exec",
            "description": "Run a command on the user's computer. Use when the user says 'run npm test', 'check disk space', 'find all TODO comments in the code', 'install the dependencies', 'start the dev server', 'check if port 3000 is in use', 'what version of node do I have'. Handles any shell command.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "The shell command to execute"},
                    "reason": {"type": "string", "description": "Brief reason why this command is needed"},
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": "Search the internet for current information. Use when the user asks about news, weather, prices, 'what's the latest version of React', 'who won the game', 'stock price of AAPL', 'temperature in Phoenix', or anything requiring up-to-date info you don't know.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "The search query"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browse_url",
            "description": "Open and read a specific webpage. Use when the user shares a URL or says 'check this article', 'read that page', 'what does this link say'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "The URL to fetch"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_schedule",
            "description": "Set a reminder or recurring task. Use when the user says 'remind me at 3pm', 'every morning check the server', 'in 30 minutes tell me to stretch', 'schedule a daily report'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "task": {"type": "string", "description": "What to do when triggered"},
                    "schedule": {"type": "string", "description": "When to run: 'every hour', 'every morning', 'in 30 minutes', 'every day at 9am'"},
                },
                "required": ["task", "schedule"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_pipeline",
            "description": "Start a multi-step project with specialized agents. Use when the user says 'build me an app', 'create a full website', 'research and write a report', 'put together a presentation', 'spawn agents to work on this'. Breaks complex work into stages with different expert agents.",
            "parameters": {
                "type": "object",
                "properties": {
                    "task": {"type": "string", "description": "Full description of what to build/research/analyze"},
                    "template": {"type": "string", "description": "Pipeline template: 'coding', 'research', 'analysis', 'drafting', 'presentation', 'general'", "default": ""},
                },
                "required": ["task"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "cancel_schedule",
            "description": "Stop a scheduled reminder or recurring task. Use when the user says 'cancel that reminder', 'stop the daily check', 'turn off the morning briefing'. Call list_schedules first to find the task ID.",
            "parameters": {
                "type": "object",
                "properties": {
                    "task_id": {"type": "string", "description": "The task ID to cancel"},
                },
                "required": ["task_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_schedules",
            "description": "Show all active reminders and scheduled tasks. Use when the user says 'what reminders do I have', 'what's scheduled', 'show my tasks', 'list my reminders'.",
            "parameters": {
                "type": "object",
                "properties": {},
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "store_memory",
            "description": "Save something to long-term memory. Use when the user says 'remember this', 'note this down', 'log this', 'save this for later', 'don't forget', 'jot this down', 'keep track of this'. Stores facts, preferences, and important information.",
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {"type": "string", "description": "What to remember"},
                },
                "required": ["content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "recall_memory",
            "description": "Search your memory for things the user told you before. Use when they say 'what did I tell you about', 'do you remember', 'what do you know about me', 'when is my meeting', 'what was that API key'. Searches all past stored facts and conversations.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "What to search for in memory"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_pptx",
            "description": "Create a PowerPoint presentation. Use when the user says 'make a presentation', 'create slides', 'put together a deck', 'make a PowerPoint about X'. Produces a real .pptx file.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path for the .pptx file (e.g. C:/Users/Jordan/Documents/report.pptx)"},
                    "title": {"type": "string", "description": "Presentation title for the title slide"},
                    "slides": {
                        "type": "array",
                        "description": "List of slides. Each slide has a 'title' and 'content' (bullet points separated by newlines).",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "content": {"type": "string", "description": "Bullet points separated by newlines"},
                            },
                        },
                    },
                },
                "required": ["path", "title", "slides"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_docx",
            "description": "Create a Word document. Use when the user says 'write a report', 'make a Word doc', 'create a document', 'draft a proposal'. Produces a real .docx file with headings and paragraphs.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path for the .docx file"},
                    "title": {"type": "string", "description": "Document title"},
                    "sections": {
                        "type": "array",
                        "description": "List of sections. Each has a 'heading' and 'body' (paragraph text).",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string"},
                                "body": {"type": "string"},
                            },
                        },
                    },
                },
                "required": ["path", "title", "sections"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_xlsx",
            "description": "Create an Excel spreadsheet. Use when the user says 'make a spreadsheet', 'create an Excel file', 'put this data in a table', 'generate a report with numbers'. Produces a real .xlsx file.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path for the .xlsx file"},
                    "sheets": {
                        "type": "array",
                        "description": "List of sheets. Each has a 'name' and 'rows' (2D array of cell values).",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string", "description": "Sheet tab name"},
                                "rows": {
                                    "type": "array",
                                    "description": "2D array: list of rows, each row is a list of cell values",
                                    "items": {"type": "array", "items": {}},
                                },
                            },
                        },
                    },
                },
                "required": ["path", "sheets"],
            },
        },
    },
]


# ---------------------------------------------------------------------------
# Role-scoped tool subsets for pipeline subagents
# ---------------------------------------------------------------------------

_ROLE_TOOL_MAP: dict[str, list[str]] = {
    "planner":      ["files_list", "files_read", "web_search", "recall_memory"],
    "backend":      ["files_list", "files_read", "files_write", "terminal_exec", "create_docx", "create_xlsx"],
    "frontend":     ["files_list", "files_read", "files_write", "terminal_exec"],
    "tester":       ["files_list", "files_read", "terminal_exec"],
    "reviewer":     ["files_list", "files_read"],
    "researcher":   ["files_list", "files_read", "web_search", "browse_url", "recall_memory"],
    "analyst":      ["files_list", "files_read", "web_search", "recall_memory", "create_xlsx"],
    "data_analyst": ["files_list", "files_read", "terminal_exec", "create_xlsx"],
    "writer":       ["files_list", "files_read", "files_write", "store_memory", "create_docx", "create_pptx"],
    "designer":     ["files_list", "files_read", "create_pptx"],
    "executor":     ["files_list", "files_read", "files_write", "terminal_exec", "create_pptx", "create_docx", "create_xlsx"],
    "presenter":    ["files_list", "files_read", "create_pptx", "create_docx"],
    "drafter":      ["files_list", "files_read", "files_write", "create_docx", "create_pptx"],
}

# Tools that pipeline subagents should NEVER get (prevents recursion, destructive ops)
_PIPELINE_BLOCKED = {"files_delete", "create_pipeline", "create_schedule", "cancel_schedule", "list_schedules"}


def get_tools_for_role(role: str) -> list[dict]:
    """Get the tool schemas allowed for a specific pipeline role."""
    allowed = _ROLE_TOOL_MAP.get(role, ["files_list", "files_read"])
    return [
        t for t in TOOL_SCHEMAS
        if t["function"]["name"] in allowed
        and t["function"]["name"] not in _PIPELINE_BLOCKED
    ]


# ---------------------------------------------------------------------------
# Tool executor — runs a tool call and returns the result as a string
# ---------------------------------------------------------------------------

def execute_tool(name: str, arguments: dict, api_port: int = 8765) -> str:
    """Execute a tool call and return the result as a string.

    This is synchronous and safe for use in both the async chat endpoint
    (via asyncio wrapping) and the synchronous pipeline handler.
    """
    base = f"http://127.0.0.1:{api_port}"
    try:
        if name == "files_list":
            dir_path = arguments.get("path", str(Path.home()))
            p = Path(dir_path)
            if not p.exists():
                return f"Directory not found: {dir_path}"
            if not p.is_dir():
                return f"Not a directory: {dir_path}"
            entries = []
            for item in sorted(p.iterdir()):
                try:
                    if item.is_dir():
                        entries.append(f"DIR  {item.name}")
                    else:
                        size = item.stat().st_size
                        entries.append(f"     {item.name} ({size} bytes)")
                except PermissionError:
                    entries.append(f"     {item.name} (access denied)")
                if len(entries) >= 50:
                    break
            if not entries:
                return f"Directory is empty: {dir_path}"
            return f"Contents of {dir_path}:\n" + "\n".join(entries)

        elif name == "files_read":
            filepath = arguments.get("path", "")
            p = Path(filepath)
            if not p.exists():
                return f"File not found: {filepath}"
            if not p.is_file():
                return f"Not a file: {filepath}"
            try:
                content = p.read_text(encoding="utf-8", errors="replace")
                lines = content.count("\n") + 1
                size = p.stat().st_size
                return f"File: {filepath} ({size} bytes, {lines} lines)\n\n{content[:3000]}"
            except Exception as e:
                return f"Error reading {filepath}: {e}"

        elif name == "files_write":
            req_path = arguments.get("path", "")
            home = str(Path.home()).replace("\\", "/")
            req_norm = req_path.replace("\\", "/")
            if not req_norm.startswith(home):
                return f"BLOCKED: Can only write files within your home directory ({home})"
            protected = [".fireside/api/", ".fireside/plugins/",
                         ".fireside/bot/", ".fireside/middleware/"]
            if any(p in req_norm for p in protected):
                return "BLOCKED: Cannot write to Fireside code directories (api/, plugins/, bot/)."
            content = arguments.get("content", "")
            if len(content) > 5_242_880:
                return f"BLOCKED: File content too large ({len(content)} bytes). Max 5MB."
            p = Path(req_path)
            try:
                p.parent.mkdir(parents=True, exist_ok=True)
                p.write_text(content, encoding="utf-8")
                return f"File saved: {req_path} ({len(content.encode('utf-8'))} bytes)"
            except Exception as e:
                return f"Error writing {req_path}: {e}"

        elif name == "files_delete":
            if not arguments.get("confirmed"):
                return (
                    f"SAFETY CHECK: Delete requires confirmation. "
                    f"TARGET: {arguments.get('path', 'unknown')}. "
                    f"Please describe exactly what will be deleted and ask the user to confirm. "
                    f"If they say yes, call files_delete again with confirmed=true."
                )
            req_path = arguments.get("path", "")
            home = str(Path.home()).replace("\\", "/")
            req_norm = req_path.replace("\\", "/")
            if not req_norm.startswith(home):
                return f"BLOCKED: Can only delete files within your home directory ({home})"
            p = Path(req_path)
            if not p.exists():
                return f"File not found: {req_path}"
            try:
                if p.is_file():
                    p.unlink()
                elif p.is_dir():
                    import shutil
                    shutil.rmtree(p)
                return f"Deleted: {req_path}"
            except Exception as e:
                return f"Error deleting {req_path}: {e}"

        elif name == "terminal_exec":
            cmd = arguments.get("command", "")
            cmd_lower = cmd.lower()

            # ALWAYS block: system-level danger regardless of path
            always_block = [
                "shutdown", "restart-computer", "stop-computer",
                "format-volume", "clear-disk", "dd if=",
                "invoke-expression", "iex ", "iex(",
                "powershell -encodedcommand", "powershell -e ",
            ]
            if any(d in cmd_lower for d in always_block):
                return (
                    f"BLOCKED: Dangerous system command: {cmd}\n"
                    f"This command is never allowed."
                )

            # PATH-AWARE check: destructive commands blocked for system dirs
            destructive_cmds = [
                "rm ", "rm -", "rmdir", "del ", "erase",
                "remove-item", "rd ", "rd/",
                "cmd /c del", "cmd /c rd",
            ]
            if any(d in cmd_lower for d in destructive_cmds):
                system_paths = [
                    "c:\\windows", "c:/windows", "/windows",
                    "c:\\program files", "c:/program files",
                    "/usr", "/bin", "/sbin", "/etc", "/boot",
                    "system32", "syswow64",
                    ".fireside\\api", ".fireside/api",
                    ".fireside\\plugins", ".fireside/plugins",
                    ".fireside\\bot", ".fireside/bot",
                ]
                root_patterns = [
                    "c:\\*", "c:/*",
                    "remove-item c:\\",
                    "rm -rf /", "rm -rf /*",
                ]
                if any(sp in cmd_lower for sp in system_paths) or any(rp in cmd_lower for rp in root_patterns):
                    return (
                        f"BLOCKED: Destructive command targets system files: {cmd}\n"
                        f"Destructive commands are only allowed within your project directories."
                    )

            try:
                result = subprocess.run(
                    cmd, shell=True, capture_output=True, text=True,
                    timeout=30, cwd=str(Path.home()),
                )
                out = f"Exit code: {result.returncode}\n"
                if result.stdout:
                    out += f"Output:\n{result.stdout[:2000]}\n"
                if result.stderr:
                    out += f"Errors:\n{result.stderr[:500]}\n"
                return out
            except subprocess.TimeoutExpired:
                return f"Command timed out after 30 seconds: {cmd}"
            except Exception as e:
                return f"Error running command: {e}"

        elif name == "web_search":
            try:
                from plugins.browse.handler import web_search as _ws
                import asyncio
                try:
                    loop = asyncio.get_running_loop()
                except RuntimeError:
                    loop = None
                if loop and loop.is_running():
                    # Already in async context — run sync wrapper
                    import concurrent.futures
                    with concurrent.futures.ThreadPoolExecutor() as pool:
                        result = pool.submit(asyncio.run, _ws(arguments.get("query", ""), max_results=5)).result()
                else:
                    result = asyncio.run(_ws(arguments.get("query", ""), max_results=5))
                if not result.get("ok"):
                    return f"Web search failed: {result.get('error', 'unknown error')}"
                if result.get("results"):
                    lines = [f"• {r['title']}: {r['snippet']}" for r in result["results"]]
                    return "\n".join(lines)
                return result.get("raw_text", "No results found.")
            except Exception as e:
                return f"Web search error: {e}"

        elif name == "browse_url":
            try:
                from plugins.browse.handler import browse as _browse
                import asyncio
                try:
                    loop = asyncio.get_running_loop()
                except RuntimeError:
                    loop = None
                if loop and loop.is_running():
                    import concurrent.futures
                    with concurrent.futures.ThreadPoolExecutor() as pool:
                        result = pool.submit(asyncio.run, _browse(arguments.get("url", ""))).result()
                else:
                    result = asyncio.run(_browse(arguments.get("url", "")))
                if not result.get("ok"):
                    return f"Failed to fetch: {result.get('error', 'unknown error')}"
                return f"Page: {result.get('title', 'Untitled')}\n\n{result.get('text', '')[:2000]}"
            except Exception as e:
                return f"Browse error: {e}"

        elif name == "create_schedule":
            payload = json.dumps({
                "task": arguments.get("task", ""),
                "schedule": arguments.get("schedule", "in 1 hour"),
                "action": "chat",
            }).encode()
            r = urllib.request.Request(f"{base}/api/v1/scheduler/create", data=payload,
                                       headers={"Content-Type": "application/json"}, method="POST")
            with urllib.request.urlopen(r, timeout=10) as resp:
                data = json.loads(resp.read())
            sched = data.get("schedule", {})
            return f"Scheduled: {sched.get('description', 'unknown')} — Task ID: {data.get('task_id', 'unknown')}"

        elif name == "create_pipeline":
            payload = json.dumps({
                "task": arguments.get("task", ""),
                "template": arguments.get("template", ""),
            }).encode()
            r = urllib.request.Request(f"{base}/api/v1/pipeline/start", data=payload,
                                       headers={"Content-Type": "application/json"}, method="POST")
            with urllib.request.urlopen(r, timeout=15) as resp:
                data = json.loads(resp.read())
            return (
                f"Pipeline created! Template: {data.get('template', 'auto')} "
                f"Status: {data.get('status', 'starting')}. "
                f"Sub-agents are working through it stage by stage."
            )

        elif name == "cancel_schedule":
            task_id = arguments.get("task_id", "")
            r = urllib.request.Request(f"{base}/api/v1/scheduler/{task_id}",
                                       method="DELETE")
            with urllib.request.urlopen(r, timeout=10) as resp:
                data = json.loads(resp.read())
            if data.get("ok"):
                return f"Schedule cancelled: {task_id}"
            return f"Failed to cancel: {data.get('error', 'unknown')}"

        elif name == "list_schedules":
            r = urllib.request.Request(f"{base}/api/v1/scheduler", method="GET")
            with urllib.request.urlopen(r, timeout=10) as resp:
                data = json.loads(resp.read())
            tasks = data.get("tasks", [])
            if not tasks:
                return "No active scheduled tasks."
            lines = []
            for t in tasks:
                sched = t.get("schedule", {}).get("description", "")
                lines.append(f"• [{t.get('id', '?')}] {t.get('task', '')[:60]} — {sched}")
            return f"Active schedules ({len(tasks)}):\n" + "\n".join(lines)

        elif name == "store_memory":
            content = arguments.get("content", "")
            try:
                import orchestrator as orch_mod
                orch_mod.observe(content, importance=0.9, source="chat_explicit")
                return f"Stored to memory: {content[:100]}"
            except Exception as ex:
                return f"Memory store failed: {ex}"

        elif name == "recall_memory":
            query = arguments.get("query", "")
            try:
                import orchestrator as orch_mod
                memories = orch_mod.recall_memories(query, top_k=5)
                if memories:
                    lines = [f"- {m.get('content', '')[:200]}" for m in memories]
                    return f"Found {len(memories)} memories:\n" + "\n".join(lines)
                return "No memories found for that topic."
            except Exception as ex:
                return f"Memory recall failed: {ex}"

        elif name == "create_pptx":
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)

                title_text = arguments.get("title", "Untitled")
                slides_data = arguments.get("slides", [])

                # Title slide
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = title_text
                if slide.placeholders[1]:
                    slide.placeholders[1].text = f"Generated by Fireside AI"

                # Content slides
                for s in slides_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = s.get("title", "")
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    content = s.get("content", "")
                    for i, line in enumerate(content.split("\n")):
                        line = line.strip()
                        if not line:
                            continue
                        if i == 0:
                            tf.text = line
                        else:
                            p = tf.add_paragraph()
                            p.text = line
                            p.level = 1 if line.startswith("-") or line.startswith("•") else 0

                out_path = arguments.get("path", "")
                p = Path(out_path)
                p.parent.mkdir(parents=True, exist_ok=True)
                prs.save(str(p))
                return f"PowerPoint created: {out_path} ({len(slides_data)} slides)"
            except ImportError:
                return "Error: python-pptx not installed. Run: pip install python-pptx"
            except Exception as e:
                return f"Error creating PowerPoint: {e}"

        elif name == "create_docx":
            try:
                from docx import Document as DocxDocument
                from docx.shared import Pt as DocxPt
                doc = DocxDocument()

                title_text = arguments.get("title", "Untitled")
                sections = arguments.get("sections", [])

                doc.add_heading(title_text, level=0)

                for sec in sections:
                    heading = sec.get("heading", "")
                    body = sec.get("body", "")
                    if heading:
                        doc.add_heading(heading, level=1)
                    if body:
                        # Split on double newlines for separate paragraphs
                        for para_text in body.split("\n\n"):
                            para_text = para_text.strip()
                            if para_text:
                                doc.add_paragraph(para_text)

                out_path = arguments.get("path", "")
                p = Path(out_path)
                p.parent.mkdir(parents=True, exist_ok=True)
                doc.save(str(p))
                return f"Word document created: {out_path} ({len(sections)} sections)"
            except ImportError:
                return "Error: python-docx not installed. Run: pip install python-docx"
            except Exception as e:
                return f"Error creating Word document: {e}"

        elif name == "create_xlsx":
            try:
                from openpyxl import Workbook
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = Workbook()
                wb.remove(wb.active)  # Remove default sheet

                sheets_data = arguments.get("sheets", [])
                if not sheets_data:
                    sheets_data = [{"name": "Sheet1", "rows": []}]

                for sheet_def in sheets_data:
                    ws = wb.create_sheet(title=sheet_def.get("name", "Sheet"))
                    rows = sheet_def.get("rows", [])

                    for r_idx, row in enumerate(rows, 1):
                        for c_idx, val in enumerate(row, 1):
                            cell = ws.cell(row=r_idx, column=c_idx, value=val)
                            # Bold the first row (header)
                            if r_idx == 1:
                                cell.font = Font(bold=True)
                                cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                                cell.font = Font(bold=True, color="FFFFFF")

                    # Auto-width columns
                    for col in ws.columns:
                        max_len = max((len(str(cell.value or "")) for cell in col), default=8)
                        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

                out_path = arguments.get("path", "")
                p = Path(out_path)
                p.parent.mkdir(parents=True, exist_ok=True)
                wb.save(str(p))
                total_rows = sum(len(s.get("rows", [])) for s in sheets_data)
                return f"Excel file created: {out_path} ({len(sheets_data)} sheets, {total_rows} rows)"
            except ImportError:
                return "Error: openpyxl not installed. Run: pip install openpyxl"
            except Exception as e:
                return f"Error creating Excel file: {e}"

        else:
            return f"Unknown tool: {name}"

    except Exception as e:
        log.warning("[tools] Tool execution failed (%s): %s", name, e)
        return f"Tool error: {str(e)}"
