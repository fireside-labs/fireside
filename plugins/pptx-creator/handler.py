"""
PowerPoint Creator — Fireside Tool Plugin.

Produces professional, polished .pptx presentations that look like
they were made by a designer, not an AI. Enterprise-quality slides.

Features:
  - Multiple professional themes (dark, light, corporate, gradient)
  - Accent bars, gradient backgrounds, proper typography
  - Chart embedding (bar, line, pie) via python-pptx charts
  - Smart layout selection based on content
  - Speaker notes support

Routes:
    POST /tools/pptx/create — Create a .pptx from slide definitions
"""

import logging
import os
import uuid
from pathlib import Path

log = logging.getLogger("valhalla.pptx")

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "output", "presentations")
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ---------------------------------------------------------------------------
# Color themes — premium palettes, not generic
# ---------------------------------------------------------------------------

THEMES = {
    "dark": {
        "bg": (0x0C, 0x0C, 0x14),
        "bg2": (0x14, 0x12, 0x1C),
        "title": (0xF0, 0xDC, 0xC8),
        "body": (0xC4, 0xA8, 0x82),
        "accent": (0xF5, 0x9E, 0x0B),
        "accent2": (0xD9, 0x77, 0x06),
        "subtle": (0x4A, 0x3D, 0x30),
        "chart_colors": [(0xF5,0x9E,0x0B), (0xA7,0x8B,0xFA), (0x34,0xD3,0x99), (0xFB,0x92,0x3C), (0x60,0xA5,0xFA)],
    },
    "light": {
        "bg": (0xFF, 0xFF, 0xFF),
        "bg2": (0xF8, 0xF8, 0xFA),
        "title": (0x1A, 0x1A, 0x2E),
        "body": (0x44, 0x44, 0x56),
        "accent": (0x25, 0x63, 0xEB),
        "accent2": (0x1D, 0x4E, 0xD8),
        "subtle": (0xA0, 0xA0, 0xB0),
        "chart_colors": [(0x25,0x63,0xEB), (0x7C,0x3A,0xED), (0x05,0x96,0x69), (0xEA,0x58,0x0C), (0xDB,0x27,0x77)],
    },
    "corporate": {
        "bg": (0x1E, 0x29, 0x3B),
        "bg2": (0x0F, 0x17, 0x2A),
        "title": (0xF1, 0xF5, 0xF9),
        "body": (0xCB, 0xD5, 0xE1),
        "accent": (0x38, 0xBD, 0xF8),
        "accent2": (0x06, 0xB6, 0xD4),
        "subtle": (0x47, 0x55, 0x69),
        "chart_colors": [(0x38,0xBD,0xF8), (0xA7,0x8B,0xFA), (0x34,0xD3,0x99), (0xFB,0xBF,0x24), (0xF4,0x72,0xB6)],
    },
    "gradient": {
        "bg": (0x0F, 0x0A, 0x1A),
        "bg2": (0x1A, 0x10, 0x30),
        "title": (0xFF, 0xFF, 0xFF),
        "body": (0xD8, 0xD0, 0xE8),
        "accent": (0xA7, 0x8B, 0xFA),
        "accent2": (0x7C, 0x3A, 0xED),
        "subtle": (0x5A, 0x4D, 0x70),
        "chart_colors": [(0xA7,0x8B,0xFA), (0xFB,0xBF,0x24), (0x34,0xD3,0x99), (0xFB,0x92,0x3C), (0xF4,0x72,0xB6)],
    },
}


def create_pptx(title: str, slides: list[dict], theme: str = "dark") -> dict:
    """
    Create a professional PowerPoint presentation.

    Each slide dict:
        - layout: "title", "section", "content", "two_column", "chart", "blank"
        - title: str
        - content: str (bullet points separated by newlines)
        - subtitle: str (for title/section slides)
        - notes: str (speaker notes)
        - chart_type: "bar", "line", "pie" (for chart layout)
        - chart_data: {labels: [...], series: [{name: str, values: [...]}]}
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed. Run: pip install python-pptx"}

    colors = THEMES.get(theme, THEMES["dark"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def rgb(t): return RGBColor(*t)
    def set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = rgb(color)

    def add_accent_bar(slide, x, y, w, h, color):
        """Add a colored accent bar."""
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        return shape

    def style_text(paragraph, text, size, color, bold=False, align=None):
        """Apply consistent text styling."""
        paragraph.text = text
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = rgb(color)
        paragraph.font.bold = bold
        paragraph.font.name = "Calibri"
        if align:
            paragraph.alignment = align

    for i, slide_def in enumerate(slides):
        layout = slide_def.get("layout", "content")
        slide_title = slide_def.get("title", f"Slide {i + 1}")
        content = slide_def.get("content", "")
        subtitle = slide_def.get("subtitle", "")
        notes = slide_def.get("notes", "")

        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        set_bg(slide, colors["bg"])

        if layout == "title":
            # ── Title Slide: centered, accent underline, subtitle ──
            add_accent_bar(slide, Inches(4.5), Inches(3.5), Inches(4.333), Emu(50000), colors["accent"])

            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            style_text(p, slide_title, 48, colors["title"], bold=True, align=PP_ALIGN.CENTER)

            if subtitle or content:
                p2 = tf.add_paragraph()
                style_text(p2, subtitle or content, 22, colors["body"], align=PP_ALIGN.CENTER)
                p2.space_before = Pt(20)

            # Bottom bar
            add_accent_bar(slide, Inches(0), Inches(7.1), Inches(13.333), Emu(60000), colors["accent2"])

        elif layout == "section":
            # ── Section divider: large text, accent stripe ──
            add_accent_bar(slide, Inches(0), Inches(0), Emu(100000), Inches(7.5), colors["accent"])

            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(3))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            style_text(p, slide_title, 44, colors["title"], bold=True)

            if subtitle or content:
                p2 = tf.add_paragraph()
                style_text(p2, subtitle or content, 18, colors["subtle"])
                p2.space_before = Pt(16)

        elif layout == "two_column":
            # ── Two column layout ──
            # Left accent bar
            add_accent_bar(slide, Inches(0.5), Inches(1.2), Emu(60000), Inches(0.6), colors["accent"])

            # Title
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(1))
            tf = txBox.text_frame
            style_text(tf.paragraphs[0], slide_title, 30, colors["accent"], bold=True)

            # Split content by \n\n for two columns, or split in half
            parts = content.split("\n\n") if "\n\n" in content else [content[:len(content)//2], content[len(content)//2:]]

            for col, (x_pos, col_content) in enumerate(zip([0.8, 7.2], parts[:2])):
                box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.8), Inches(5.5), Inches(5))
                tf = box.text_frame
                tf.word_wrap = True
                bullets = [b.strip() for b in col_content.split("\n") if b.strip()]
                for j, bullet in enumerate(bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f"▸ {bullet}" if not bullet.startswith(("▸", "•", "-", "→")) else bullet
                    p.font.size = Pt(16)
                    p.font.color.rgb = rgb(colors["body"])
                    p.font.name = "Calibri"
                    p.space_after = Pt(8)

        elif layout == "chart":
            # ── Chart slide ──
            add_accent_bar(slide, Inches(0.5), Inches(1.0), Emu(60000), Inches(0.5), colors["accent"])

            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(12), Inches(1))
            tf = txBox.text_frame
            style_text(tf.paragraphs[0], slide_title, 28, colors["accent"], bold=True)

            chart_type = slide_def.get("chart_type", "bar")
            chart_data = slide_def.get("chart_data", {})

            if chart_data:
                try:
                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE

                    data = CategoryChartData()
                    data.categories = chart_data.get("labels", ["A", "B", "C"])
                    for series in chart_data.get("series", [{"name": "Data", "values": [1, 2, 3]}]):
                        data.add_series(series["name"], series["values"])

                    chart_types = {
                        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                        "line": XL_CHART_TYPE.LINE_MARKERS,
                        "pie": XL_CHART_TYPE.PIE,
                    }
                    ct = chart_types.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
                    chart_frame = slide.shapes.add_chart(
                        ct, Inches(1), Inches(1.5), Inches(11.333), Inches(5.5), data
                    )

                    # Style the chart
                    chart = chart_frame.chart
                    chart.has_legend = len(chart_data.get("series", [])) > 1
                    if chart.has_legend:
                        chart.legend.include_in_layout = False

                    # Color the series
                    for idx, series in enumerate(chart.series):
                        fill = series.format.fill
                        fill.solid()
                        c = colors["chart_colors"][idx % len(colors["chart_colors"])]
                        fill.fore_color.rgb = rgb(c)

                except Exception as e:
                    # Fallback: show data as text
                    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
                    tf = box.text_frame
                    tf.word_wrap = True
                    style_text(tf.paragraphs[0], f"Chart data: {chart_data}", 14, colors["body"])

        elif layout == "blank":
            pass

        else:
            # ── Default content slide: accent bar + bullets ──
            # Left accent bar
            add_accent_bar(slide, Inches(0.5), Inches(1.2), Emu(60000), Inches(0.6), colors["accent"])

            # Title
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(1))
            tf = txBox.text_frame
            style_text(tf.paragraphs[0], slide_title, 30, colors["accent"], bold=True)

            # Horizontal rule under title
            add_accent_bar(slide, Inches(0.8), Inches(1.35), Inches(2.5), Emu(25000), colors["accent2"])

            # Body content
            bullets = [b.strip() for b in content.split("\n") if b.strip()]
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.2))
            tf = body_box.text_frame
            tf.word_wrap = True

            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()

                # Smart formatting: detect sub-bullets vs main bullets
                if bullet.startswith(("  ", "\t")):
                    p.text = f"    ◦ {bullet.strip()}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = rgb(colors["subtle"])
                else:
                    clean = bullet.lstrip("-•▸→ ").strip()
                    p.text = f"▸  {clean}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = rgb(colors["body"])

                p.font.name = "Calibri"
                p.space_after = Pt(10)

            # Slide number (bottom right)
            num_box = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.8), Inches(0.4))
            style_text(num_box.text_frame.paragraphs[0], str(i + 1), 10, colors["subtle"], align=PP_ALIGN.RIGHT)

        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # Save
    filename = f"{title.lower().replace(' ', '_')[:40]}_{uuid.uuid4().hex[:6]}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)

    log.info("[pptx] Created %s with %d slides (theme: %s)", filename, len(slides), theme)

    return {
        "ok": True,
        "filename": filename,
        "filepath": filepath,
        "slide_count": len(slides),
        "theme": theme,
    }


# ---------------------------------------------------------------------------
# DOCX / PPTX → PDF conversion
# ---------------------------------------------------------------------------

def convert_to_pdf(input_path: str) -> dict:
    """
    Convert a DOCX or PPTX file to PDF.

    Tries in order:
      1. docx2pdf (Windows — uses Word/PowerPoint COM automation)
      2. LibreOffice headless (cross-platform fallback)
    """
    p = Path(input_path)
    if not p.exists():
        return {"ok": False, "error": f"File not found: {input_path}"}

    ext = p.suffix.lower()
    if ext not in (".docx", ".pptx", ".doc", ".ppt", ".xlsx"):
        return {"ok": False, "error": f"Unsupported format: {ext}. Supports: docx, pptx, xlsx"}

    pdf_path = str(p.with_suffix(".pdf"))

    # Method 1: docx2pdf (Windows, requires Office installed)
    try:
        import docx2pdf
        docx2pdf.convert(str(p), pdf_path)
        return {
            "ok": True,
            "pdf_path": pdf_path,
            "method": "docx2pdf (Microsoft Office)",
        }
    except ImportError:
        pass
    except Exception as e:
        log.warning("[pptx] docx2pdf failed: %s", e)

    # Method 2: LibreOffice headless
    import subprocess, shutil
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if lo:
        try:
            output_dir = str(p.parent)
            result = subprocess.run(
                [lo, "--headless", "--convert-to", "pdf", "--outdir", output_dir, str(p)],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0 and Path(pdf_path).exists():
                return {
                    "ok": True,
                    "pdf_path": pdf_path,
                    "method": "LibreOffice",
                }
        except Exception as e:
            log.warning("[pptx] LibreOffice conversion failed: %s", e)

    return {
        "ok": False,
        "error": "PDF conversion requires either: pip install docx2pdf (+ MS Office) or LibreOffice installed.",
    }


# ---------------------------------------------------------------------------
# Template Learning — extract styles from existing .pptx files
# ---------------------------------------------------------------------------

TEMPLATE_DIR = Path.home() / ".valhalla" / "templates"
TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)


def _rgb_to_tuple(rgb_color) -> tuple:
    """Convert pptx RGBColor to (R, G, B) tuple."""
    try:
        return (rgb_color[0], rgb_color[1], rgb_color[2])
    except Exception:
        return (0x44, 0x44, 0x44)


def scan_template(pptx_path: str) -> dict:
    """
    Scan an existing .pptx and extract its visual style:
      - Background colors
      - Title and body font families + sizes
      - Accent colors (from shapes, lines)
      - Layout patterns (what type of content on each slide)
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches, Emu
        from pptx.dml.color import RGBColor
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed"}

    p = Path(pptx_path)
    if not p.exists():
        return {"ok": False, "error": f"File not found: {pptx_path}"}

    prs = Presentation(str(p))

    # Collect data from all slides
    bg_colors = []
    title_fonts = []
    title_sizes = []
    body_fonts = []
    body_sizes = []
    accent_colors = []
    layout_sequence = []
    shape_colors = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {"index": slide_idx, "shapes": 0, "has_chart": False,
                      "has_image": False, "text_blocks": 0}

        # Background
        try:
            bg = slide.background
            if bg.fill and bg.fill.fore_color and bg.fill.fore_color.rgb:
                c = bg.fill.fore_color.rgb
                bg_colors.append((c[0], c[1], c[2]))
        except Exception:
            pass

        for shape in slide.shapes:
            slide_info["shapes"] += 1

            # Check shape types
            if shape.has_chart:
                slide_info["has_chart"] = True
            if hasattr(shape, 'image'):
                slide_info["has_image"] = True

            # Shape fill colors (for accent detection)
            try:
                if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
                    c = shape.fill.fore_color.rgb
                    shape_colors.append((c[0], c[1], c[2]))
            except Exception:
                pass

            # Text styling
            if shape.has_text_frame:
                slide_info["text_blocks"] += 1
                for para in shape.text_frame.paragraphs:
                    if not para.text.strip():
                        continue
                    try:
                        font = para.runs[0].font if para.runs else para.font
                        fname = font.name or "Calibri"
                        fsize = font.size

                        # Heuristic: large text = title, smaller = body
                        if fsize and fsize >= Pt(24):
                            title_fonts.append(fname)
                            title_sizes.append(fsize.pt)
                        elif fsize:
                            body_fonts.append(fname)
                            body_sizes.append(fsize.pt)
                        else:
                            body_fonts.append(fname)

                        # Text colors
                        if font.color and font.color.rgb:
                            c = font.color.rgb
                            color_tuple = (c[0], c[1], c[2])
                            # Light colors on dark bg = title, etc.
                            brightness = sum(color_tuple) / 3
                            if brightness > 200:
                                pass  # White/near-white text — probably title on dark
                            elif brightness > 100:
                                accent_colors.append(color_tuple)
                    except Exception:
                        pass

        # Determine layout pattern
        if slide_info["has_chart"]:
            layout_sequence.append("chart")
        elif slide_info["has_image"]:
            layout_sequence.append("image")
        elif slide_info["text_blocks"] <= 1 and slide_idx == 0:
            layout_sequence.append("title")
        elif slide_info["shapes"] > 5:
            layout_sequence.append("complex")
        else:
            layout_sequence.append("content")

    # Determine most common values
    def most_common(lst, default):
        if not lst:
            return default
        from collections import Counter
        return Counter(lst).most_common(1)[0][0]

    def avg(lst, default):
        if not lst:
            return default
        return round(sum(lst) / len(lst))

    # Filter accent colors: remove near-black and near-white
    accent_colors = [c for c in accent_colors + shape_colors
                     if 30 < sum(c)/3 < 230]

    # Build the template profile
    profile = {
        "source_file": str(p.name),
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "colors": {
            "bg": most_common(bg_colors, (0x0C, 0x0C, 0x14)),
            "title": most_common(
                [(c[0], c[1], c[2]) for c in accent_colors if sum(c)/3 > 150],
                (0xF0, 0xDC, 0xC8)
            ),
            "body": most_common(
                [(c[0], c[1], c[2]) for c in accent_colors if 80 < sum(c)/3 < 180],
                (0xC4, 0xA8, 0x82)
            ),
            "accent": most_common(accent_colors, (0xF5, 0x9E, 0x0B)),
            "accent2": accent_colors[1] if len(accent_colors) > 1 else (0xD9, 0x77, 0x06),
            "subtle": (0x4A, 0x3D, 0x30),
        },
        "fonts": {
            "title_font": most_common(title_fonts, "Calibri"),
            "title_size": avg(title_sizes, 36),
            "body_font": most_common(body_fonts, "Calibri"),
            "body_size": avg(body_sizes, 16),
        },
        "layout_sequence": layout_sequence,
    }

    # Serialize slide dimensions
    profile["slide_width"] = int(profile["slide_width"])
    profile["slide_height"] = int(profile["slide_height"])

    return {"ok": True, "profile": profile}


def learn_template(pptx_path: str, template_name: str = "") -> dict:
    """
    Scan a .pptx and save the extracted style as a reusable template.
    User can then say "use my company template" and the AI will apply it.
    """
    result = scan_template(pptx_path)
    if not result.get("ok"):
        return result

    profile = result["profile"]

    # Auto-name from source file if not provided
    if not template_name:
        template_name = Path(pptx_path).stem.replace(" ", "_").lower()

    # Save template
    template_path = TEMPLATE_DIR / f"{template_name}.json"
    template_path.write_text(json.dumps(profile, indent=2), encoding="utf-8")

    # Also store in long-term memory for ambient awareness
    try:
        import orchestrator as orch_mod
        orch_mod.observe(
            f"[Template learned: '{template_name}'] "
            f"Colors: bg={profile['colors']['bg']}, accent={profile['colors']['accent']}. "
            f"Fonts: {profile['fonts']['title_font']} for titles, {profile['fonts']['body_font']} for body. "
            f"Source: {profile['source_file']}",
            importance=0.8,
            source="template_learning",
        )
    except Exception:
        pass

    log.info("[pptx] Learned template '%s' from %s", template_name, pptx_path)

    return {
        "ok": True,
        "template_name": template_name,
        "template_path": str(template_path),
        "profile": profile,
    }


def list_templates() -> dict:
    """List all saved templates."""
    templates = []
    for f in sorted(TEMPLATE_DIR.glob("*.json")):
        try:
            data = json.loads(f.read_text(encoding="utf-8"))
            templates.append({
                "name": f.stem,
                "source": data.get("source_file", "unknown"),
                "title_font": data.get("fonts", {}).get("title_font", "Calibri"),
                "accent": data.get("colors", {}).get("accent", []),
            })
        except Exception:
            pass
    return {"ok": True, "templates": templates, "count": len(templates)}


def _load_template_as_theme(template_name: str) -> dict | None:
    """Load a saved template and convert to theme format usable by create_pptx."""
    path = TEMPLATE_DIR / f"{template_name}.json"
    if not path.exists():
        return None
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        colors = data.get("colors", {})
        # Convert to THEMES format
        bg = tuple(colors.get("bg", (0x0C, 0x0C, 0x14)))
        accent = tuple(colors.get("accent", (0xF5, 0x9E, 0x0B)))
        accent2 = tuple(colors.get("accent2", (0xD9, 0x77, 0x06)))
        theme = {
            "bg": bg,
            "bg2": tuple(max(0, c - 10) for c in bg),
            "title": tuple(colors.get("title", (0xF0, 0xDC, 0xC8))),
            "body": tuple(colors.get("body", (0xC4, 0xA8, 0x82))),
            "accent": accent,
            "accent2": accent2,
            "subtle": tuple(colors.get("subtle", (0x4A, 0x3D, 0x30))),
            "chart_colors": [accent, accent2, (0x34, 0xD3, 0x99), (0xFB, 0x92, 0x3C), (0x60, 0xA5, 0xFA)],
            # Include font info
            "_fonts": data.get("fonts", {}),
        }
        return theme
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Smart auto-layout detection
# ---------------------------------------------------------------------------

def auto_detect_layout(slide_def: dict, slide_index: int, total_slides: int) -> str:
    """
    Automatically pick the best layout based on content.
    The LLM can still override by specifying layout explicitly.
    """
    content = slide_def.get("content", "")
    title = slide_def.get("title", "")
    chart_data = slide_def.get("chart_data", {})

    # First slide = title
    if slide_index == 0 and not content:
        return "title"

    # Has chart data → chart layout
    if chart_data:
        return "chart"

    # Two distinct sections separated by double newline → two column
    if "\n\n" in content and len(content.split("\n\n")) == 2:
        return "two_column"

    # Very short content (section divider)
    bullets = [b for b in content.split("\n") if b.strip()]
    if len(bullets) <= 1 and len(content) < 50:
        return "section"

    # Default
    return "content"


# ---------------------------------------------------------------------------
# FastAPI route registration
# ---------------------------------------------------------------------------

import json

def register_routes(app, config: dict = None):
    """Register PowerPoint creation + template learning routes."""
    from pydantic import BaseModel

    class SlideSpec(BaseModel):
        layout: str = "content"
        title: str = ""
        content: str = ""
        subtitle: str = ""
        notes: str = ""
        chart_type: str = ""
        chart_data: dict = {}

    class PptxRequest(BaseModel):
        title: str
        slides: list[SlideSpec]
        theme: str = "dark"
        template: str = ""

    class ConvertRequest(BaseModel):
        path: str

    class LearnRequest(BaseModel):
        path: str
        name: str = ""

    @app.post("/tools/pptx/create")
    async def handle_create(req: PptxRequest):
        slides = [s.model_dump() for s in req.slides]
        theme = req.theme
        # If template specified, load it as a theme
        if req.template:
            learned = _load_template_as_theme(req.template)
            if learned:
                THEMES[f"_learned_{req.template}"] = learned
                theme = f"_learned_{req.template}"
        return create_pptx(req.title, slides, theme)

    @app.post("/tools/convert-to-pdf")
    async def handle_convert(req: ConvertRequest):
        return convert_to_pdf(req.path)

    @app.post("/tools/pptx/learn-template")
    async def handle_learn(req: LearnRequest):
        return learn_template(req.path, req.name)

    @app.get("/tools/pptx/templates")
    async def handle_list_templates():
        return list_templates()

    @app.post("/tools/pptx/scan-template")
    async def handle_scan(req: LearnRequest):
        return scan_template(req.path)

    log.info("[pptx] Routes registered (with template learning)")

